"""
Document processing utilities for extracting text from various file formats.
Supports PDF, PPTX, and DOCX files.
"""
import os
from typing import Dict, List, Optional
from pathlib import Path
import logging

import PyPDF2
import pdfplumber
from pptx import Presentation
from docx import Document

logger = logging.getLogger(__name__)


class DocumentProcessor:
    """Process and extract text from various document formats."""
    
    SUPPORTED_FORMATS = ['.pdf', '.pptx', '.docx', '.txt']
    
    def __init__(self):
        """Initialize the document processor."""
        pass
    
    def process_file(self, file_path: str) -> Dict[str, any]:
        """
        Process a file and extract its text content.
        
        Args:
            file_path: Path to the file to process
            
        Returns:
            Dictionary containing extracted text and metadata
        """
        file_path = Path(file_path)
        
        if not file_path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")
        
        file_ext = file_path.suffix.lower()
        
        if file_ext not in self.SUPPORTED_FORMATS:
            raise ValueError(f"Unsupported file format: {file_ext}")
        
        try:
            if file_ext == '.pdf':
                text, metadata = self._process_pdf(file_path)
            elif file_ext == '.pptx':
                text, metadata = self._process_pptx(file_path)
            elif file_ext == '.docx':
                text, metadata = self._process_docx(file_path)
            elif file_ext == '.txt':
                text, metadata = self._process_txt(file_path)
            else:
                raise ValueError(f"Unsupported file format: {file_ext}")
            
            return {
                'text': text,
                'metadata': metadata,
                'file_path': str(file_path),
                'file_name': file_path.name,
                'file_type': file_ext
            }
        except Exception as e:
            logger.error(f"Error processing file {file_path}: {str(e)}")
            raise
    
    def _process_pdf(self, file_path: Path) -> tuple[str, Dict]:
        """
        Extract text from PDF file using pdfplumber.
        
        Args:
            file_path: Path to the PDF file
            
        Returns:
            Tuple of (extracted_text, metadata)
        """
        text_content = []
        metadata = {
            'page_count': 0,
            'extraction_method': 'pdfplumber'
        }
        
        try:
            with pdfplumber.open(file_path) as pdf:
                metadata['page_count'] = len(pdf.pages)
                
                for page_num, page in enumerate(pdf.pages, 1):
                    page_text = page.extract_text()
                    if page_text:
                        text_content.append(f"\n--- Page {page_num} ---\n")
                        text_content.append(page_text)
                
                # Extract PDF metadata if available
                if pdf.metadata:
                    metadata.update({
                        'title': pdf.metadata.get('Title', ''),
                        'author': pdf.metadata.get('Author', ''),
                        'subject': pdf.metadata.get('Subject', ''),
                        'creator': pdf.metadata.get('Creator', '')
                    })
        except Exception as e:
            logger.warning(f"pdfplumber failed for {file_path}, trying PyPDF2: {str(e)}")
            # Fallback to PyPDF2
            text_content, metadata = self._process_pdf_pypdf2(file_path)
        
        return '\n'.join(text_content), metadata
    
    def _process_pdf_pypdf2(self, file_path: Path) -> tuple[str, Dict]:
        """
        Fallback method to extract text from PDF using PyPDF2.
        
        Args:
            file_path: Path to the PDF file
            
        Returns:
            Tuple of (extracted_text, metadata)
        """
        text_content = []
        metadata = {
            'page_count': 0,
            'extraction_method': 'PyPDF2'
        }
        
        with open(file_path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            metadata['page_count'] = len(pdf_reader.pages)
            
            for page_num, page in enumerate(pdf_reader.pages, 1):
                page_text = page.extract_text()
                if page_text:
                    text_content.append(f"\n--- Page {page_num} ---\n")
                    text_content.append(page_text)
            
            # Extract PDF metadata
            if pdf_reader.metadata:
                metadata.update({
                    'title': pdf_reader.metadata.get('/Title', ''),
                    'author': pdf_reader.metadata.get('/Author', ''),
                    'subject': pdf_reader.metadata.get('/Subject', ''),
                    'creator': pdf_reader.metadata.get('/Creator', '')
                })
        
        return '\n'.join(text_content), metadata
    
    def _process_pptx(self, file_path: Path) -> tuple[str, Dict]:
        """
        Extract text from PowerPoint file.
        
        Args:
            file_path: Path to the PPTX file
            
        Returns:
            Tuple of (extracted_text, metadata)
        """
        text_content = []
        metadata = {
            'slide_count': 0,
            'extraction_method': 'python-pptx'
        }
        
        prs = Presentation(file_path)
        metadata['slide_count'] = len(prs.slides)
        
        for slide_num, slide in enumerate(prs.slides, 1):
            text_content.append(f"\n--- Slide {slide_num} ---\n")
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text_content.append(shape.text)
                
                # Extract text from tables
                if shape.has_table:
                    for row in shape.table.rows:
                        row_text = ' | '.join(cell.text for cell in row.cells)
                        text_content.append(row_text)
        
        return '\n'.join(text_content), metadata
    
    def _process_docx(self, file_path: Path) -> tuple[str, Dict]:
        """
        Extract text from Word document.
        
        Args:
            file_path: Path to the DOCX file
            
        Returns:
            Tuple of (extracted_text, metadata)
        """
        text_content = []
        metadata = {
            'paragraph_count': 0,
            'extraction_method': 'python-docx'
        }
        
        doc = Document(file_path)
        
        # Extract paragraphs
        for para in doc.paragraphs:
            if para.text.strip():
                text_content.append(para.text)
        
        metadata['paragraph_count'] = len(doc.paragraphs)
        
        # Extract text from tables
        for table in doc.tables:
            for row in table.rows:
                row_text = ' | '.join(cell.text for cell in row.cells)
                text_content.append(row_text)
        
        # Extract core properties
        core_props = doc.core_properties
        metadata.update({
            'title': core_props.title or '',
            'author': core_props.author or '',
            'subject': core_props.subject or ''
        })
        
        return '\n'.join(text_content), metadata
    
    def _process_txt(self, file_path: Path) -> tuple[str, Dict]:
        """
        Extract text from plain text file.
        
        Args:
            file_path: Path to the TXT file
            
        Returns:
            Tuple of (extracted_text, metadata)
        """
        with open(file_path, 'r', encoding='utf-8') as file:
            text = file.read()
        
        metadata = {
            'char_count': len(text),
            'line_count': text.count('\n') + 1,
            'extraction_method': 'plain_text'
        }
        
        return text, metadata


class TextChunker:
    """Split text into chunks for embedding and retrieval using token-based chunking."""
    
    def __init__(self, chunk_size: int = 500, chunk_overlap: int = 100):
        """
        Initialize the text chunker with TOKEN-based chunking.
        
        Args:
            chunk_size: Maximum size of each chunk in TOKENS (not characters)
            chunk_overlap: Number of overlapping TOKENS between chunks
        """
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap
        
        # Validate inputs
        if chunk_size <= 0:
            raise ValueError(f"chunk_size must be positive, got {chunk_size}")
        
        if chunk_overlap < 0:
            raise ValueError(f"chunk_overlap must be non-negative, got {chunk_overlap}")
        
        if chunk_overlap >= chunk_size:
            logger.warning(
                f"  chunk_overlap ({chunk_overlap}) >= chunk_size ({chunk_size}). "
                f"This may cause excessive overlapping chunks. "
                f"Recommended: overlap < 50% of chunk_size"
            )
        
        # Initialize tiktoken for token-based chunking
        try:
            import tiktoken
            self.encoding = tiktoken.get_encoding("cl100k_base")  # Used by text-embedding-3-* models
            self.use_tokens = True
            logger.info(f" Using TOKEN-based chunking with tiktoken (limit: {chunk_size} tokens)")
        except ImportError:
            logger.warning("  tiktoken not installed. Falling back to CHARACTER-based chunking (not recommended)")
            logger.warning("Install tiktoken: pip install tiktoken")
            self.encoding = None
            self.use_tokens = False
    
    def chunk_text(self, text: str, metadata: Optional[Dict] = None) -> List[Dict]:
        """
        Split text into overlapping chunks using TOKEN-based chunking (if available).
        
        Args:
            text: Text to chunk
            metadata: Optional metadata to attach to each chunk
            
        Returns:
            List of dictionaries containing chunk text and metadata
        """
        if self.use_tokens:
            return self._chunk_text_by_tokens(text, metadata)
        else:
            return self._chunk_text_by_characters(text, metadata)
    
    def _chunk_text_by_tokens(self, text: str, metadata: Optional[Dict] = None) -> List[Dict]:
        """
        Split text into chunks based on TOKEN count (recommended for OpenAI).
        
        Args:
            text: Text to chunk
            metadata: Optional metadata to attach to each chunk
            
        Returns:
            List of dictionaries containing chunk text and metadata
        """
        logger.info(f"Starting TOKEN-based chunking - text length: {len(text)} characters")
        
        if not text or not text.strip():
            logger.warning("Empty text provided for chunking")
            return []
        
        # Encode entire text to tokens
        tokens = self.encoding.encode(text)
        total_tokens = len(tokens)
        logger.info(f"Total tokens: {total_tokens} (chunk_size: {self.chunk_size} tokens, overlap: {self.chunk_overlap} tokens)")
        
        chunks = []
        start_token_idx = 0
        previous_start = -1  # Track previous position to detect infinite loops
        
        while start_token_idx < total_tokens:
            # Prevent infinite loop: ensure we're making forward progress
            if start_token_idx == previous_start:
                logger.error(f"Infinite loop detected at token {start_token_idx}. Breaking.")
                break
            previous_start = start_token_idx
            
            # Calculate end token index
            end_token_idx = min(start_token_idx + self.chunk_size, total_tokens)
            
            # Extract tokens for this chunk
            chunk_tokens = tokens[start_token_idx:end_token_idx]
            
            # Decode tokens back to text
            chunk_text = self.encoding.decode(chunk_tokens).strip()
            
            # Only add non-empty chunks (skip chunks that are too small to be meaningful)
            if chunk_text and len(chunk_text) >= 3:  # At least 3 characters for meaningful content
                chunk_data = {
                    'text': chunk_text,
                    'token_count': len(chunk_tokens),
                    'chunk_index': len(chunks)
                }
                
                if metadata:
                    chunk_data['metadata'] = metadata.copy()
                
                chunks.append(chunk_data)
            
            # If we've reached the end of the text, break
            if end_token_idx >= total_tokens:
                break
            
            # Move to next chunk with overlap
            next_start = end_token_idx - self.chunk_overlap
            
            # Prevent moving backwards or staying in same position
            # If overlap is too large, just move forward by 1 token minimum
            if next_start <= start_token_idx:
                next_start = start_token_idx + max(1, self.chunk_size // 2)
                logger.warning(f"Overlap ({self.chunk_overlap}) too large for chunk_size ({self.chunk_size}). "
                             f"Adjusting to prevent infinite loop.")
            
            start_token_idx = next_start
            
            # Log progress every 50 chunks
            if len(chunks) % 50 == 0 and len(chunks) > 0:
                progress = (start_token_idx / total_tokens) * 100
                logger.info(f"  Chunking progress: {progress:.1f}% ({len(chunks)} chunks created)")
        
        logger.info(f" TOKEN-based chunking completed - created {len(chunks)} chunks from {total_tokens} tokens")
        return chunks
    
    def _chunk_text_by_characters(self, text: str, metadata: Optional[Dict] = None) -> List[Dict]:
        """
        Split text into chunks based on CHARACTER count (fallback method).
        
        Args:
            text: Text to chunk
            metadata: Optional metadata to attach to each chunk
            
        Returns:
            List of dictionaries containing chunk text and metadata
        """
        logger.info(f"Starting CHARACTER-based chunking (fallback) - text length: {len(text)} characters")
        
        if not text or not text.strip():
            logger.warning("Empty text provided for chunking")
            return []
        
        chunks = []
        start = 0
        text_length = len(text)
        previous_start = -1  # Track previous position to detect infinite loops
        logger.info(f"Chunk size: {self.chunk_size}, overlap: {self.chunk_overlap}")
        
        while start < text_length:
            # Prevent infinite loop: ensure we're making forward progress
            if start == previous_start:
                logger.error(f"Infinite loop detected at position {start}. Breaking.")
                break
            previous_start = start
            end = start + self.chunk_size
            
            # If this is not the last chunk, try to break at a sentence or word boundary
            if end < text_length:
                # Look for sentence boundary (., !, ?)
                for punct in ['. ', '! ', '? ', '\n\n']:
                    last_punct = text.rfind(punct, start, end)
                    if last_punct != -1:
                        end = last_punct + 1
                        break
                else:
                    # If no sentence boundary, look for word boundary
                    last_space = text.rfind(' ', start, end)
                    if last_space != -1:
                        end = last_space
            
            chunk_text = text[start:end].strip()
            
            if chunk_text:
                chunk_data = {
                    'text': chunk_text,
                    'start_index': start,
                    'end_index': end,
                    'chunk_index': len(chunks)
                }
                
                if metadata:
                    chunk_data['metadata'] = metadata.copy()
                
                chunks.append(chunk_data)
            
            # If we've reached the end, break
            if end >= text_length:
                break
            
            # Move to next chunk with overlap
            next_start = end - self.chunk_overlap
            
            # Prevent moving backwards or staying in same position
            # If overlap is too large, just move forward by at least half chunk size
            if next_start <= start:
                next_start = start + max(1, self.chunk_size // 2)
                logger.warning(f"Overlap ({self.chunk_overlap}) too large for chunk_size ({self.chunk_size}). "
                             f"Adjusting to prevent infinite loop.")
            
            start = next_start
        
        logger.info(f" CHARACTER-based chunking completed - created {len(chunks)} chunks")
        return chunks
